"""External document converters and binary discovery."""

import os
import re
import shutil
import subprocess
import tempfile
import threading
from pathlib import Path

LOCAL_APPDATA = Path(os.environ.get("LOCALAPPDATA", Path.home() / "AppData" / "Local"))

PANDOC_CANDIDATES = [
    shutil.which("pandoc"),
    *LOCAL_APPDATA.glob("Microsoft/WinGet/Packages/JohnMacFarlane.Pandoc_*/pandoc-*/pandoc.exe"),
]
PANDOC_BIN = next((p for p in PANDOC_CANDIDATES if p and Path(p).exists()), None)

UNHWP_CANDIDATES = [
    shutil.which("unhwp"),
    LOCAL_APPDATA / "Microsoft" / "WindowsApps" / "unhwp.exe",
]
UNHWP_BIN = next((p for p in UNHWP_CANDIDATES if p and Path(p).exists()), None)

if os.name == "nt":
    TESSDATA_DIR = LOCAL_APPDATA / "tessdata"
else:
    tessdata_candidates = [
        Path(os.environ["TESSDATA_PREFIX"]) if os.environ.get("TESSDATA_PREFIX") else None,
        Path("/opt/homebrew/share/tessdata"),
        Path("/usr/local/share/tessdata"),
        Path("/usr/share/tesseract-ocr/5/tessdata"),
    ]
    TESSDATA_DIR = next(
        (path for path in tessdata_candidates if path and path.exists()),
        Path("/usr/local/share/tessdata"),
    )

TESSERACT_BIN = shutil.which("tesseract") or str(
    Path(os.environ.get("ProgramFiles", r"C:\Program Files")) / "Tesseract-OCR" / "tesseract.exe"
)
POPPLER_BIN = shutil.which("pdftoppm") or next(
    (
        str(path)
        for path in LOCAL_APPDATA.glob(
            "Microsoft/WinGet/Packages/oschwartz10612.Poppler_*/poppler-*/Library/bin/pdftoppm.exe"
        )
    ),
    None,
)

PDF_CPU_LIMIT = 4
PDF_TIMEOUT_SECONDS = 1800
PAGE_IMAGE_RE = re.compile(r"-(\d+)\.png$")


def convert_with_unhwp(src: Path) -> str:
    """Convert HWP/HWPX to Markdown with unhwp."""
    if not UNHWP_BIN:
        raise RuntimeError("unhwp was not found")
    with tempfile.TemporaryDirectory() as td:
        local_src = Path(td) / src.name
        output_md = Path(td) / "extract.md"
        shutil.copy2(src, local_src)
        result = subprocess.run(
            [str(UNHWP_BIN), "md", local_src.name, "-o", output_md.name],
            cwd=td,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
        )
        if result.returncode != 0 or not output_md.exists():
            raise RuntimeError(
                f"unhwp did not produce Markdown (exit {result.returncode}): "
                f"{result.stdout.strip()} {result.stderr.strip()}"
            )
        return output_md.read_text(encoding="utf-8")


def convert_with_markitdown(src: Path) -> str:
    """Convert a document with Microsoft MarkItDown."""
    from markitdown import MarkItDown

    try:
        return MarkItDown().convert(str(src)).text_content
    except (KeyboardInterrupt, SystemExit):
        raise
    except BaseException as exc:
        # Some MarkItDown conversion errors inherit directly from BaseException.
        raise RuntimeError(str(exc)) from exc


def convert_with_pandoc(src: Path) -> str:
    """Convert a document with Pandoc."""
    if not PANDOC_BIN:
        raise RuntimeError("pandoc was not found")
    result = subprocess.run(
        [str(PANDOC_BIN), str(src), "-t", "markdown", "--wrap=none"],
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        check=True,
    )
    return result.stdout


def convert_with_pymupdf(src: Path) -> str:
    """Extract plain PDF text with PyMuPDF."""
    import fitz

    doc = fitz.open(str(src))
    try:
        return "".join(page.get_text() for page in doc)
    finally:
        doc.close()


def _limit_pdf_cpu() -> None:
    """Reduce CPU pressure where psutil exposes affinity controls."""
    try:
        import psutil

        process = psutil.Process()
        if hasattr(process, "cpu_affinity"):
            available = process.cpu_affinity()
            if available:
                process.cpu_affinity(available[: min(PDF_CPU_LIMIT, len(available))])
    except Exception:
        pass


def convert_with_pymupdf4llm(src: Path, timeout: float = PDF_TIMEOUT_SECONDS) -> list[str]:
    """Extract structured Markdown from a PDF while preserving page units."""
    _limit_pdf_cpu()
    outcome: dict[str, object] = {}

    def run() -> None:
        import pymupdf4llm

        try:
            outcome["pages"] = pymupdf4llm.to_markdown(str(src), page_chunks=True)
        except Exception as exc:
            outcome["error"] = exc

    worker = threading.Thread(target=run, daemon=True)
    worker.start()
    worker.join(timeout)
    if worker.is_alive():
        raise TimeoutError(f"pymupdf4llm exceeded {timeout}s on {src.name}")
    if "error" in outcome:
        raise outcome["error"]  # type: ignore[misc]
    return [page["text"] for page in outcome["pages"]]  # type: ignore[index]


def convert_with_ocr(src: Path, pages: list[int] | None = None):
    """OCR a whole PDF or selected zero-based pages with Tesseract kor+eng."""
    if not POPPLER_BIN:
        raise RuntimeError("pdftoppm was not found")
    if not Path(TESSERACT_BIN).exists():
        raise RuntimeError("tesseract was not found")

    with tempfile.TemporaryDirectory() as td:
        prefix = Path(td) / "page"
        if pages is None:
            subprocess.run(
                [str(POPPLER_BIN), "-png", "-r", "200", str(src), str(prefix)],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                check=False,
            )
        else:
            for start, end in _contiguous_ranges(sorted(pages)):
                subprocess.run(
                    [
                        str(POPPLER_BIN),
                        "-png",
                        "-r",
                        "200",
                        "-f",
                        str(start + 1),
                        "-l",
                        str(end + 1),
                        str(src),
                        str(prefix),
                    ],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    check=False,
                )

        page_images = sorted(
            Path(td).glob("page-*.png"),
            key=lambda p: int(PAGE_IMAGE_RE.search(p.name).group(1)),
        )
        if not page_images:
            raise RuntimeError("pdftoppm produced no page images")

        chunks: list[str] = []
        by_page: dict[int, str] = {}
        out_base = Path(td) / "page_out"
        for png_path in page_images:
            subprocess.run(
                [
                    TESSERACT_BIN,
                    str(png_path),
                    str(out_base),
                    "-l",
                    "kor+eng",
                    "--tessdata-dir",
                    str(TESSDATA_DIR),
                ],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                check=False,
            )
            out_txt = out_base.with_suffix(".txt")
            text = out_txt.read_text(encoding="utf-8", errors="replace") if out_txt.exists() else ""
            if out_txt.exists():
                out_txt.unlink()

            page_number = int(PAGE_IMAGE_RE.search(png_path.name).group(1))
            png_path.unlink()
            if pages is None:
                chunks.append(text)
            else:
                by_page[page_number - 1] = text

        return by_page if pages is not None else "\n\n".join(chunks)


def _shape_markdown(shape, mso_shape_type) -> list[str]:
    """Extract text/table/chart content from a PowerPoint shape recursively."""
    if shape.shape_type == mso_shape_type.GROUP:
        parts: list[str] = []
        for child in shape.shapes:
            parts.extend(_shape_markdown(child, mso_shape_type))
        return parts

    parts: list[str] = []
    if shape.has_table:
        parts.append(_table_markdown(shape.table))
    if shape.has_chart:
        try:
            chart_text = _chart_markdown(shape.chart)
        except Exception:
            chart_text = ""
        if chart_text:
            parts.append(chart_text)
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            parts.append(text)
    return parts


def _table_markdown(table) -> str:
    rows = [
        [cell.text.strip().replace("\n", " ").replace("|", "\\|") for cell in row.cells]
        for row in table.rows
    ]
    if not rows:
        return ""
    lines = ["| " + " | ".join(row) + " |" for row in rows]
    lines.insert(1, "|" + "|".join(["---"] * len(rows[0])) + "|")
    return "\n".join(lines)


def _chart_markdown(chart) -> str:
    if not chart.has_title:
        return ""
    title = chart.chart_title.text_frame.text.strip()
    return f"**{title}**" if title else ""


def convert_with_pptx(src: Path) -> list[str]:
    """Extract PowerPoint content while preserving slide boundaries."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(src))
    slides: list[str] = []
    for slide in prs.slides:
        title_shape = slide.shapes.title
        parts: list[str] = []
        if title_shape is not None and title_shape.has_text_frame:
            title_text = title_shape.text_frame.text.strip()
            if title_text:
                parts.append(f"# {title_text}")
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            parts.extend(_shape_markdown(shape, MSO_SHAPE_TYPE))
        slides.append("\n\n".join(parts))
    return slides


def _contiguous_ranges(indexes: list[int]) -> list[tuple[int, int]]:
    ranges: list[tuple[int, int]] = []
    start = prev = None
    for index in indexes:
        if start is None:
            start = prev = index
        elif index == prev + 1:
            prev = index
        else:
            ranges.append((start, prev))
            start = prev = index
    if start is not None:
        ranges.append((start, prev))
    return ranges
